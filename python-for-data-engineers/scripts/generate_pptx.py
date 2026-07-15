#!/usr/bin/env python3
"""
Generate PowerPoint (.pptx) decks for all 12 sessions.

Usage:
    python scripts/generate_pptx.py
    python scripts/generate_pptx.py --session 1

Requires: pip install python-pptx
"""

from __future__ import annotations

import argparse
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
MARP_DIR = ROOT / "slides" / "marp"
OUT_DIR = ROOT / "slides" / "pptx"

# Brand colors
NAVY = RGBColor(0x1A, 0x52, 0x76)
BLUE = RGBColor(0x28, 0x74, 0xA6)
DARK = RGBColor(0x2C, 0x3E, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)


def parse_marp_slides(md_path: Path) -> list[dict]:
    """Parse Marp markdown into slide dicts with title and bullets."""
    text = md_path.read_text(encoding="utf-8")
    # Strip front matter
    if text.startswith("---"):
        parts = text.split("---", 2)
        text = parts[2] if len(parts) > 2 else text

    raw_slides = re.split(r"\n---\n", text.strip())
    slides = []

    for raw in raw_slides:
        lines = raw.strip().splitlines()
        if not lines:
            continue

        title = ""
        subtitle = ""
        bullets: list[str] = []
        code_blocks: list[str] = []
        in_code = False
        code_buf: list[str] = []

        for line in lines:
            if line.strip().startswith("```"):
                if in_code:
                    code_blocks.append("\n".join(code_buf))
                    code_buf = []
                    in_code = False
                else:
                    in_code = True
                continue
            if in_code:
                code_buf.append(line)
                continue
            if line.startswith("# ") and not title:
                title = line[2:].strip()
            elif line.startswith("## ") and not subtitle:
                subtitle = line[3:].strip()
            elif line.startswith("- ") or line.startswith("* "):
                bullets.append(line[2:].strip())
            elif line.startswith("|") and "---" not in line:
                bullets.append(line.strip())

        if not title and lines:
            title = lines[0].lstrip("#").strip()

        slides.append(
            {
                "title": title,
                "subtitle": subtitle,
                "bullets": bullets,
                "code": code_blocks,
            }
        )
    return slides


def _set_title_style(shape, size: int = 36, color: RGBColor = NAVY) -> None:
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        para.font.size = Pt(size)
        para.font.bold = True
        para.font.color.rgb = color


def add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    _set_title_style(slide.shapes.title, 40, NAVY)
    if len(slide.placeholders) > 1:
        for para in slide.placeholders[1].text_frame.paragraphs:
            para.font.size = Pt(22)
            para.font.color.rgb = BLUE


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    code: list[str] | None = None,
) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    _set_title_style(slide.shapes.title, 32, NAVY)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, bullet in enumerate(bullets[:8]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = DARK

    if code:
        top = Inches(4.2) if bullets else Inches(1.8)
        box = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(2.8))
        tf_code = box.text_frame
        tf_code.word_wrap = True
        p = tf_code.paragraphs[0]
        p.text = code[0][:1200]
        p.font.name = "Courier New"
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        box.fill.solid()
        box.fill.fore_color.rgb = CODE_BG


def build_pptx(session_num: int, md_path: Path, out_path: Path) -> None:
    slides = parse_marp_slides(md_path)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for i, s in enumerate(slides):
        if i == 0:
            sub = s.get("subtitle") or f"Session {session_num:02d}"
            add_title_slide(prs, s["title"], sub)
        else:
            add_content_slide(prs, s["title"] or s.get("subtitle", "Slide"), s["bullets"], s.get("code"))

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Created {out_path} ({len(slides)} slides)")


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate PPTX from Marp slides")
    parser.add_argument("--session", type=int, help="Single session number (1-12)")
    args = parser.parse_args()

    sessions = [args.session] if args.session else range(1, 13)
    for n in sessions:
        md = MARP_DIR / f"session-{n:02d}.md"
        if not md.exists():
            print(f"Skip missing {md}")
            continue
        out = OUT_DIR / f"session-{n:02d}.pptx"
        build_pptx(n, md, out)


if __name__ == "__main__":
    main()
